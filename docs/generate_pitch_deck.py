"""Generate the SafeRide pitch deck (.pptx) for the Africa Ignite Hackathon.

Run:
    python docs/generate_pitch_deck.py

Output: docs/SafeRide_Pitch_Deck.pptx (also FR/EN content embedded).

Sections (per HackerEarth Africa Ignite evaluation criteria):
1. Cover
2. Problem statement & context
3. Proposed solution & API usage
4. Technical architecture
5. Business model & monetization
6. Team members & roles
7. Business value summary
8. API usage summary & commercial impact
9. Closing / Call to action
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt


# ----- Brand palette --------------------------------------------------------
BG_DARK = RGBColor(0x06, 0x09, 0x0F)
BG_PANEL = RGBColor(0x0B, 0x14, 0x24)
ACCENT = RGBColor(0x22, 0xD3, 0xEE)         # cyan
ACCENT_2 = RGBColor(0x22, 0xC5, 0x5E)       # green
WARN = RGBColor(0xFB, 0xBF, 0x24)
DANGER = RGBColor(0xF8, 0x71, 0x71)
TEXT = RGBColor(0xE6, 0xEC, 0xF7)
MUTED = RGBColor(0x94, 0xA3, 0xB8)


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _set_bg(slide, color=BG_DARK):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    return bg


def _accent_bar(slide, top=Inches(0.4), height=Inches(0.08), width=Inches(1.6), left=Inches(0.6)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    return bar


def _add_text(slide, text, left, top, width, height, *, size=18, bold=False,
              color=TEXT, align=PP_ALIGN.LEFT, font="Inter"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def _bullets(slide, bullets, left, top, width, height, *, size=16, color=TEXT,
             gap=6, bold_first=False, bullet_color=ACCENT_2):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        run = p.add_run()
        run.text = "▸  " + item
        run.font.name = "Inter"
        run.font.size = Pt(size)
        if bold_first and i == 0:
            run.font.bold = True
        run.font.color.rgb = color
    return tb


def _panel(slide, left, top, width, height, fill=BG_PANEL, accent=False):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    rect.adjustments[0] = 0.06
    rect.line.color.rgb = ACCENT if accent else RGBColor(0x1F, 0x2A, 0x44)
    rect.line.width = Pt(0.75)
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    rect.shadow.inherit = False
    return rect


def _kpi(slide, value, label, left, top, width=Inches(2.4), height=Inches(1.4)):
    _panel(slide, left, top, width, height)
    _add_text(slide, value, left, top + Inches(0.18), width, Inches(0.7),
              size=30, bold=True, color=ACCENT, align=PP_ALIGN.CENTER, font="Space Grotesk")
    _add_text(slide, label, left, top + Inches(0.85), width, Inches(0.4),
              size=12, color=MUTED, align=PP_ALIGN.CENTER)


def _section_title(slide, eyebrow, title):
    _accent_bar(slide)
    _add_text(slide, eyebrow, Inches(0.6), Inches(0.55), Inches(8), Inches(0.4),
              size=12, bold=True, color=ACCENT)
    _add_text(slide, title, Inches(0.6), Inches(0.85), Inches(12.2), Inches(0.9),
              size=32, bold=True, color=TEXT, font="Space Grotesk")


def _footer(slide, idx, total):
    _add_text(slide, "SafeRide  ·  Africa Ignite Hackathon  ·  Nokia Network-as-Code",
              Inches(0.6), Inches(7.05), Inches(10), Inches(0.3),
              size=10, color=MUTED)
    _add_text(slide, f"{idx} / {total}",
              Inches(11.6), Inches(7.05), Inches(1.2), Inches(0.3),
              size=10, color=MUTED, align=PP_ALIGN.RIGHT)


# =============================================================================
# Slide builders
# =============================================================================


def slide_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s)

    # Decorative gradient blobs (faked with semi-transparent ellipses)
    for left, top, w, h, color in [
        (Inches(8.5), Inches(-1), Inches(7), Inches(7), ACCENT),
        (Inches(-2), Inches(4.5), Inches(6), Inches(6), ACCENT_2),
    ]:
        blob = s.shapes.add_shape(MSO_SHAPE.OVAL, left, top, w, h)
        blob.line.fill.background()
        blob.fill.solid()
        blob.fill.fore_color.rgb = color
        blob.fill.fore_color.brightness = -0.5

    _accent_bar(s, top=Inches(2.8), left=Inches(0.9), width=Inches(2.0))
    _add_text(s, "SafeRide", Inches(0.9), Inches(3.0), Inches(11), Inches(1.3),
              size=88, bold=True, color=TEXT, font="Space Grotesk")
    _add_text(s, "Every ride, checked. Every journey, safe.",
              Inches(0.9), Inches(4.4), Inches(11), Inches(0.7),
              size=24, color=ACCENT)
    _add_text(s, "Network Trust Score for Sub-Saharan ride-hailing",
              Inches(0.9), Inches(5.1), Inches(11), Inches(0.6),
              size=18, color=MUTED)
    _add_text(s, "Powered by 7 CAMARA APIs via Nokia Network-as-Code  ·  Africa Ignite Hackathon 2026",
              Inches(0.9), Inches(6.6), Inches(11), Inches(0.4),
              size=12, color=MUTED)


def slide_problem(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s)
    _section_title(s, "01  ·  Problem & Context", "Trust is broken in African ride-hailing")

    _bullets(s, [
        "Sub-Saharan ride-hailing market: USD 14B by 2030 (GSMA).",
        "Fake drivers using SIM-swapped accounts → assaults & thefts.",
        "Operational fraud: spoofed GPS to cheat allocation engines.",
        "Service discontinuity: GPS / voice cut off mid-ride.",
        "Existing apps (Bolt, Uber) only rely on subjective ratings.",
    ], Inches(0.6), Inches(1.95), Inches(7.4), Inches(4.5), size=18, gap=10)

    # KPI panel right
    _panel(s, Inches(8.4), Inches(1.95), Inches(4.4), Inches(4.6))
    _add_text(s, "Key data (GSMA 2024 / CAMARA)", Inches(8.6), Inches(2.05),
              Inches(4.0), Inches(0.4), size=14, bold=True, color=ACCENT)
    _kpi(s, "+340%", "SIM swap incidents (Sub-Saharan, 2023)",
         Inches(8.6), Inches(2.55), Inches(4.0), Inches(1.0))
    _kpi(s, "1.2 B$", "Mobile fraud losses (2023)",
         Inches(8.6), Inches(3.65), Inches(4.0), Inches(1.0))
    _kpi(s, "42%", "Trust in ride-hailing apps",
         Inches(8.6), Inches(4.75), Inches(4.0), Inches(1.0))

    _footer(s, 2, 9)


def slide_solution(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s)
    _section_title(s, "02  ·  Proposed Solution & API Usage",
                   "A Network Trust Score before & during every ride")

    _add_text(s, "SafeRide orchestrates 7 CAMARA APIs via Nokia Network-as-Code to compute a real-time score (0–100).",
              Inches(0.6), Inches(1.85), Inches(12.0), Inches(0.5),
              size=16, color=MUTED)

    formula = "SCR  =  moteur propriétaire SafeRide combinant signaux réseau, contexte et historique"
    _panel(s, Inches(0.6), Inches(2.45), Inches(12.0), Inches(0.85), accent=True)
    _add_text(s, formula, Inches(0.6), Inches(2.6), Inches(12.0), Inches(0.6),
              size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER, font="Space Grotesk")

    # Three trust buckets
    buckets = [
        ("70 – 100", "Reliable", "Ride allowed, standard tracking", ACCENT_2),
        ("40 – 69",  "Attention", "Offered with passenger warning", WARN),
        ("0 – 39",   "Blocked",   "Driver excluded · SafeRide alert", DANGER),
    ]
    for i, (rng, label, desc, color) in enumerate(buckets):
        left = Inches(0.6 + i * 4.15)
        _panel(s, left, Inches(3.55), Inches(4.0), Inches(1.65))
        _add_text(s, rng, left + Inches(0.2), Inches(3.7), Inches(3.6), Inches(0.5),
                  size=22, bold=True, color=color, font="Space Grotesk")
        _add_text(s, label, left + Inches(0.2), Inches(4.18), Inches(3.6), Inches(0.4),
                  size=16, bold=True, color=TEXT)
        _add_text(s, desc, left + Inches(0.2), Inches(4.55), Inches(3.6), Inches(0.6),
                  size=12, color=MUTED)

    # APIs grid
    apis = [
        ("SIM Swap", "Anti-Fraud · MANDATORY"),
        ("Number Verification", "Anti-Fraud · MANDATORY"),
        ("Location Verification", "Anti-Fraud · MANDATORY"),
        ("Device Swap", "Anti-Fraud · MANDATORY"),
        ("Quality on Demand", "Connectivity · CORE"),
        ("Congestion Insights", "Network Intelligence"),
        ("Geofencing", "Network Intelligence"),
    ]
    for i, (name, cat) in enumerate(apis):
        left = Inches(0.6 + (i % 4) * 3.15)
        top = Inches(5.45 + (i // 4) * 0.85)
        _panel(s, left, top, Inches(3.0), Inches(0.75))
        _add_text(s, name, left + Inches(0.15), top + Inches(0.05),
                  Inches(2.7), Inches(0.35), size=12, bold=True, color=TEXT)
        _add_text(s, cat, left + Inches(0.15), top + Inches(0.4),
                  Inches(2.7), Inches(0.3), size=10, color=ACCENT)

    _footer(s, 3, 9)


def slide_architecture(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s)
    _section_title(s, "03  ·  Technical Architecture",
                   "Cloud-native, operator-agnostic, AI-orchestrated")

    layers = [
        ("Mobile / Web Frontend", "React 18 + Vite + Leaflet + i18n FR/EN", ACCENT),
        ("API Gateway", "FastAPI (Python) · OpenAPI auto · CORS", ACCENT),
        ("Trust Agent (Agentic AI)", "Order APIs by region · explainable log · LLM-ready", ACCENT_2),
        ("Nokia Network-as-Code SDK", "7 CAMARA APIs orchestrated in parallel", ACCENT_2),
        ("Operator Network (MTN, Orange, Safaricom...)", "Real-time, tamper-proof network signals", WARN),
    ]

    top = Inches(2.05)
    for i, (title, desc, color) in enumerate(layers):
        y = top + Inches(i * 0.95)
        _panel(s, Inches(0.6), y, Inches(8.5), Inches(0.85))
        _add_text(s, title, Inches(0.85), y + Inches(0.08), Inches(8.0), Inches(0.4),
                  size=15, bold=True, color=color)
        _add_text(s, desc, Inches(0.85), y + Inches(0.45), Inches(8.0), Inches(0.4),
                  size=12, color=MUTED)
        if i < len(layers) - 1:
            _add_text(s, "↓", Inches(4.6), y + Inches(0.85),
                      Inches(0.4), Inches(0.2), size=14, color=MUTED, align=PP_ALIGN.CENTER)

    # Side stack
    _panel(s, Inches(9.4), Inches(2.05), Inches(3.5), Inches(4.7))
    _add_text(s, "Data plane", Inches(9.6), Inches(2.15), Inches(3.2), Inches(0.4),
              size=14, bold=True, color=ACCENT)
    _bullets(s, [
        "PostgreSQL (drivers, rides, subs)",
        "Redis (real-time sessions)",
        "Firebase Cloud Messaging",
        "OpenStreetMap + Leaflet",
        "Chariow checkout (iframe + webhook)",
        "Bearer-secured admin API",
        "OpenTelemetry + Sentry (Phase 1)",
    ], Inches(9.6), Inches(2.6), Inches(3.2), Inches(4.0), size=12, gap=6)

    _footer(s, 4, 9)


def slide_business_model(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s)
    _section_title(s, "04  ·  Business Model & Monetization",
                   "Four revenue streams, aligned with operators")

    streams = [
        ("Ride commission", "8%", "vs 15–25% Bolt/Uber", ACCENT_2),
        ("Premium driver", "5,000 FCFA / month", "via Chariow (Mobile Money, cards, crypto)", ACCENT),
        ("B2B Trust Score API", "Per-call licensing", "Logistics, fintech, delivery, healthcare", WARN),
        ("Operator partnership", "Revenue share", "Per CAMARA API call (Open Gateway model)", ACCENT_2),
    ]
    for i, (name, value, sub, color) in enumerate(streams):
        col, row = i % 2, i // 2
        left = Inches(0.6 + col * 6.3)
        top = Inches(2.0 + row * 2.1)
        _panel(s, left, top, Inches(6.0), Inches(1.85))
        _add_text(s, name, left + Inches(0.3), top + Inches(0.18),
                  Inches(5.5), Inches(0.45), size=16, bold=True, color=TEXT)
        _add_text(s, value, left + Inches(0.3), top + Inches(0.65),
                  Inches(5.5), Inches(0.55), size=24, bold=True, color=color, font="Space Grotesk")
        _add_text(s, sub, left + Inches(0.3), top + Inches(1.25),
                  Inches(5.5), Inches(0.5), size=12, color=MUTED)

    _add_text(s, "Phase 2 projection · 10,000 active users  →  ≈ 29 MFCFA / month MRR",
              Inches(0.6), Inches(6.45), Inches(12.2), Inches(0.5),
              size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    _footer(s, 5, 9)


def slide_team(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s)
    _section_title(s, "05  ·  Team Members & Roles",
                   "A multidisciplinary team for telco-grade mobility")

    team = [
        ("Jules Ndanga", "Founder & Tech Lead",
         "Backend FastAPI, Trust Agent orchestration, Nokia NaC integration."),
        ("Product / UX Lead", "Passenger & Driver experience",
         "Bilingual flows (FR / EN), accessibility, mobile-first design."),
        ("Telecom Partnership Lead", "Operator relationships",
         "MTN Cameroon, Orange, Safaricom · CAMARA compliance · GSMA Open Gateway."),
        ("Data / AI Engineer", "Trust scoring & fraud patterns",
         "LLM agent (Llama / Mistral), regional fraud signature learning, RAG."),
        ("Operations & Compliance", "Local rollout",
         "City managers (Douala, Yaoundé), data privacy (CDP, NDPR, DPA)."),
    ]
    for i, (name, role, desc) in enumerate(team):
        col, row = i % 3, i // 3
        left = Inches(0.6 + col * 4.2)
        top = Inches(2.0 + row * 2.4)
        _panel(s, left, top, Inches(3.95), Inches(2.15))
        # Avatar circle
        av = s.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.25),
                                top + Inches(0.2), Inches(0.6), Inches(0.6))
        av.line.fill.background()
        av.fill.solid()
        av.fill.fore_color.rgb = ACCENT if i % 2 == 0 else ACCENT_2
        _add_text(s, name.split()[0][0] + (name.split()[1][0] if len(name.split()) > 1 else ""),
                  left + Inches(0.25), top + Inches(0.27), Inches(0.6), Inches(0.45),
                  size=18, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER, font="Space Grotesk")
        _add_text(s, name, left + Inches(0.95), top + Inches(0.2),
                  Inches(2.85), Inches(0.4), size=15, bold=True, color=TEXT)
        _add_text(s, role, left + Inches(0.95), top + Inches(0.55),
                  Inches(2.85), Inches(0.4), size=11, color=ACCENT)
        _add_text(s, desc, left + Inches(0.25), top + Inches(0.95),
                  Inches(3.55), Inches(1.1), size=11, color=MUTED)

    _footer(s, 6, 9)


def slide_business_value(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s)
    _section_title(s, "06  ·  Business Value Summary",
                   "Why SafeRide creates measurable value, fast")

    cards = [
        ("Passengers", "Reduced incidents, transparent trust, instant disclosure of risk."),
        ("Drivers", "Network-verified reputation, priority allocation when Premium."),
        ("Mobile operators", "New revenue per CAMARA API call, demonstrate Open Gateway value."),
        ("Tech ecosystem", "Reusable Trust Score for fintech, logistics, delivery, healthcare."),
        ("Regulators", "Tamper-proof KYC for mobility — measurable fraud reduction."),
        ("12-month target", "100,000 active users · fraud rate < 0.5% · NPS > 70."),
    ]
    for i, (k, v) in enumerate(cards):
        col, row = i % 3, i // 3
        left = Inches(0.6 + col * 4.2)
        top = Inches(2.0 + row * 2.3)
        _panel(s, left, top, Inches(3.95), Inches(2.05))
        _add_text(s, k, left + Inches(0.25), top + Inches(0.2),
                  Inches(3.55), Inches(0.45), size=15, bold=True, color=ACCENT)
        _add_text(s, v, left + Inches(0.25), top + Inches(0.7),
                  Inches(3.55), Inches(1.3), size=12, color=TEXT)

    _footer(s, 7, 9)


def slide_api_impact(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s)
    _section_title(s, "07  ·  API Usage Summary & Commercial Impact",
                   "7 CAMARA APIs, one operational outcome")

    rows = [
        ("SIM Swap",            "Detect SIM exchange < 72h", "Eliminates account-takeover fraud"),
        ("Number Verification", "Bind phone ↔ carrier",      "Strong KYC without paperwork"),
        ("Device Swap",         "Detect compromised device",  "Stops handset-based impersonation"),
        ("Location Verification","Cross-check GPS vs network","Kills GPS-spoofing fraud"),
        ("Quality on Demand",   "Guarantee 4G/5G during ride","Reliable real-time tracking"),
        ("Congestion Insights", "Anticipate network drops",   "Proactive rerouting"),
        ("Geofencing",          "Alert on perimeter exit",    "Live safety alerting"),
    ]

    # header
    _panel(s, Inches(0.6), Inches(1.95), Inches(12.2), Inches(0.55), accent=True)
    headers = [("CAMARA API", 0.85), ("Network signal", 4.85), ("Commercial impact", 8.85)]
    for label, x in headers:
        _add_text(s, label, Inches(x), Inches(2.05), Inches(4.0), Inches(0.4),
                  size=13, bold=True, color=ACCENT)

    for i, (api, signal, impact) in enumerate(rows):
        y = Inches(2.6 + i * 0.55)
        _panel(s, Inches(0.6), y, Inches(12.2), Inches(0.5))
        _add_text(s, api, Inches(0.85), y + Inches(0.08), Inches(4.0), Inches(0.4),
                  size=13, bold=True, color=TEXT)
        _add_text(s, signal, Inches(4.85), y + Inches(0.08), Inches(4.0), Inches(0.4),
                  size=12, color=MUTED)
        _add_text(s, impact, Inches(8.85), y + Inches(0.08), Inches(4.0), Inches(0.4),
                  size=12, color=ACCENT_2)

    _add_text(s, "Source: https://networkascode.nokia.io/products/digital-identity-and-anti-fraud",
              Inches(0.6), Inches(6.55), Inches(12.0), Inches(0.4),
              size=11, color=MUTED)

    _footer(s, 8, 9)


def slide_close(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s)
    _section_title(s, "08  ·  Roadmap & Ask",
                   "From hackathon to 5 cities in 18 months")

    phases = [
        ("Phase 0", "Apr 24 – May 10, 2026", "Hackathon prototype · 7 CAMARA APIs simulated", ACCENT),
        ("Phase 1", "Jun – Aug 2026",        "Douala pilot · 500 drivers · MTN partnership · real NaC", ACCENT_2),
        ("Phase 2", "Sep – Dec 2026",        "Yaoundé + Lagos · 10,000 users · Premium launch", WARN),
        ("Phase 3", "2027",                  "Dakar · Nairobi · Abidjan · B2B Trust Score license", DANGER),
    ]
    for i, (phase, period, desc, color) in enumerate(phases):
        y = Inches(2.0 + i * 0.95)
        _panel(s, Inches(0.6), y, Inches(8.5), Inches(0.85))
        _add_text(s, phase, Inches(0.85), y + Inches(0.1), Inches(1.5), Inches(0.4),
                  size=14, bold=True, color=color, font="Space Grotesk")
        _add_text(s, period, Inches(2.4), y + Inches(0.1), Inches(2.4), Inches(0.4),
                  size=12, bold=True, color=TEXT)
        _add_text(s, desc, Inches(0.85), y + Inches(0.45), Inches(8.0), Inches(0.4),
                  size=12, color=MUTED)

    _panel(s, Inches(9.4), Inches(2.0), Inches(3.5), Inches(4.65), accent=True)
    _add_text(s, "Funding ask", Inches(9.6), Inches(2.15), Inches(3.2), Inches(0.4),
              size=14, bold=True, color=ACCENT)
    _add_text(s, "120 – 180 k€", Inches(9.6), Inches(2.6), Inches(3.2), Inches(0.7),
              size=28, bold=True, color=TEXT, font="Space Grotesk")
    _add_text(s, "9 months · Douala pilot", Inches(9.6), Inches(3.3), Inches(3.2), Inches(0.4),
              size=12, color=MUTED)
    _bullets(s, [
        "Cloud + observability",
        "Operator integrations",
        "Driver acquisition",
        "Compliance (CDP, RGPD)",
        "City managers x 3",
    ], Inches(9.6), Inches(3.85), Inches(3.2), Inches(2.6), size=12, gap=6)

    _add_text(s, "SafeRide doesn’t just move passengers — it turns every ride into proof of network identity.",
              Inches(0.6), Inches(6.7), Inches(12.2), Inches(0.5),
              size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    _footer(s, 9, 9)


def main() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_cover(prs)
    slide_problem(prs)
    slide_solution(prs)
    slide_architecture(prs)
    slide_business_model(prs)
    slide_team(prs)
    slide_business_value(prs)
    slide_api_impact(prs)
    slide_close(prs)

    out = Path(__file__).parent / "SafeRide_Pitch_Deck.pptx"
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    main()
